import subprocess
import sys
from pathlib import Path
import pytest

SCRIPT = Path(__file__).parents[1] / "skills/aem/edge-delivery-services/aem-doc-converter/scripts/md-to-pptx.py"


def run_script(input_path, output_path, template=None):
    cmd = [sys.executable, str(SCRIPT), str(input_path), str(output_path)]
    if template:
        cmd += ["--template", str(template)]
    return subprocess.run(cmd, capture_output=True, text=True)


def test_produces_pptx_file(training_md_path, tmp_path):
    out = tmp_path / "out.pptx"
    result = run_script(training_md_path, out)
    assert result.returncode == 0
    assert out.exists()
    assert out.stat().st_size > 0


def test_h1_becomes_title_slide(training_md_path, tmp_path):
    from pptx import Presentation
    out = tmp_path / "out.pptx"
    run_script(training_md_path, out)
    prs = Presentation(str(out))
    assert len(prs.slides) > 0
    first_slide_text = " ".join(
        shape.text for shape in prs.slides[0].shapes if shape.has_text_frame
    )
    assert "Exercise 1" in first_slide_text or "Cards Block" in first_slide_text


def test_h2_sections_become_slides(training_md_path, tmp_path):
    from pptx import Presentation
    out = tmp_path / "out.pptx"
    run_script(training_md_path, out)
    prs = Presentation(str(out))
    # training.md has H1 + H2 Prerequisites + H2 Background + H2 Steps + H2 Verification + H2 References
    assert len(prs.slides) >= 4


def test_missing_input_exits_nonzero(tmp_path):
    out = tmp_path / "out.pptx"
    result = run_script(tmp_path / "nonexistent.md", out)
    assert result.returncode != 0


def test_nonexistent_template_falls_back_to_default(training_md_path, tmp_path):
    out = tmp_path / "out.pptx"
    result = run_script(training_md_path, out, template=tmp_path / "fake-template.pptx")
    assert result.returncode == 0
    assert out.exists()
