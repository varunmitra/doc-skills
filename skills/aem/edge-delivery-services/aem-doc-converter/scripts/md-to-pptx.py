#!/usr/bin/env python3
"""
md-to-pptx.py — Convert structured Markdown to PowerPoint.

Mapping:
  H1 → Title slide (large title + subtitle from next paragraph)
  H2 → Section header slide
  H3 → Content slide title (bullets below become content)
  Bullet lists → Content bullets on current slide
  Code blocks → Code text box on current slide
  --- → Slide break

Usage: python md-to-pptx.py <input.md> <output.pptx> [--template <template.pptx>]
"""
import sys
import re
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: Missing dependency. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _get_layout(prs, preferred_idx, fallback_idx=1):
    """Get slide layout by index, falling back gracefully."""
    try:
        return prs.slide_layouts[preferred_idx]
    except IndexError:
        return prs.slide_layouts[min(fallback_idx, len(prs.slide_layouts) - 1)]


def make_title_slide(prs, title_text, subtitle_text=""):
    slide = prs.slides.add_slide(_get_layout(prs, 0))
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1 and subtitle_text:
        try:
            slide.placeholders[1].text = subtitle_text
        except (KeyError, IndexError):
            pass
    return slide


def make_section_slide(prs, title_text):
    slide = prs.slides.add_slide(_get_layout(prs, 2, fallback_idx=1))
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    return slide


def make_content_slide(prs, title_text, bullets, code_blocks):
    slide = prs.slides.add_slide(_get_layout(prs, 1))
    if slide.shapes.title:
        slide.shapes.title.text = title_text

    if bullets:
        content_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                content_ph = ph
                break
        if content_ph:
            tf = content_ph.text_frame
            tf.clear()
            for idx, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0

    if code_blocks:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "\n".join(code_blocks)
        run.font.name = "Courier New"
        run.font.size = Pt(9)
    return slide


def _strip_inline_md(text):
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    return text.strip()


def parse_slides(source):
    """Parse markdown source into a list of slide descriptors."""
    lines = source.splitlines()
    slides = []
    current = None
    in_code = False
    code_buf = []

    def flush():
        if current:
            slides.append(dict(current))

    for line in lines:
        stripped = line.strip()

        if stripped.startswith("```"):
            if not in_code:
                in_code = True
                code_buf = []
            else:
                in_code = False
                if current:
                    current.setdefault("code", []).append("\n".join(code_buf))
                code_buf = []
            continue

        if in_code:
            code_buf.append(line)
            continue

        h1 = re.match(r'^#\s+(.*)', stripped)
        h2 = re.match(r'^##\s+(.*)', stripped)
        h3 = re.match(r'^###\s+(.*)', stripped)
        bullet = re.match(r'^[-*]\s+(.*)', stripped)
        num = re.match(r'^\d+\.\s+(.*)', stripped)
        hr = re.match(r'^---+$', stripped)

        if h1:
            flush()
            current = {"type": "title", "title": _strip_inline_md(h1.group(1)),
                       "subtitle": "", "bullets": [], "code": []}
        elif h2:
            flush()
            current = {"type": "section", "title": _strip_inline_md(h2.group(1)),
                       "bullets": [], "code": []}
        elif h3:
            flush()
            current = {"type": "content", "title": _strip_inline_md(h3.group(1)),
                       "bullets": [], "code": []}
        elif bullet and current:
            current.setdefault("bullets", []).append(_strip_inline_md(bullet.group(1)))
        elif num and current:
            current.setdefault("bullets", []).append(_strip_inline_md(num.group(1)))
        elif hr:
            flush()
            current = None
        elif stripped and current and current.get("type") == "title" and not current.get("subtitle"):
            current["subtitle"] = _strip_inline_md(stripped)

    flush()
    return slides


def convert(input_path, output_path, template_path=None):
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
    else:
        if template_path:
            print(f"Warning: Template not found at {template_path}, using default layout.",
                  file=sys.stderr)
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    source = Path(input_path).read_text(encoding="utf-8")
    slides = parse_slides(source)

    for slide_data in slides:
        stype = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        if stype == "title":
            make_title_slide(prs, title, slide_data.get("subtitle", ""))
        elif stype == "section":
            make_section_slide(prs, title)
        else:
            make_content_slide(prs, title,
                               slide_data.get("bullets", []),
                               slide_data.get("code", []))

    if not prs.slides:
        make_title_slide(prs, Path(input_path).stem)

    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Convert Markdown to PPTX")
    parser.add_argument("input", help="Input .md file")
    parser.add_argument("output", help="Output .pptx file")
    parser.add_argument("--template", help="Adobe brand .pptx template", default=None)
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    convert(input_path, args.output, args.template)
    print(f"✓ Converted to {args.output}")


if __name__ == "__main__":
    main()
